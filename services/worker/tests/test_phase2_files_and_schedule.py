from __future__ import annotations

import hashlib
import io
import os
import tarfile
import uuid
import zipfile
from datetime import UTC, datetime
from pathlib import Path

os.environ.update(
    {
        "APP_ENV": "test",
        "APP_MODE": "demo",
        "DATABASE_URL": "sqlite+pysqlite:///:memory:",
        "SESSION_SECRET": "worker-test-session-secret-at-least-32-chars",
        "CSRF_SECRET": "worker-test-csrf-secret-at-least-32-chars",
        "AUDIT_HMAC_KEY": "worker-test-audit-key-at-least-32-characters",
    }
)

from docx import Document
from openpyxl import Workbook
from pptx import Presentation
from pypdf import PdfWriter
from pypdf.generic import DecodedStreamObject, DictionaryObject, NameObject

from executive_ai_worker.assistant_orchestrator import _query_terms, _reciprocal_rank_fusion
from executive_ai_worker.embedding_cache import EmbeddingArtifact, preload_artifact
from executive_ai_worker.file_extraction import (
    ExtractedBlock,
    FileExtractionPermanentError,
    _embedding_model,
    chunk_blocks,
    parse_file,
)
from executive_ai_worker.scheduler import next_cron_time


def _docx() -> bytes:
    document = Document()
    document.add_heading("经营复盘", level=1)
    document.add_paragraph("华东事业部回款进度需要关注。")
    output = io.BytesIO()
    document.save(output)
    return output.getvalue()


def _xlsx() -> bytes:
    workbook = Workbook()
    worksheet = workbook.active
    worksheet.title = "回款"
    worksheet.append(["客户", "金额"])
    worksheet.append(["演示客户001", 4200000])
    output = io.BytesIO()
    workbook.save(output)
    return output.getvalue()


def _pptx() -> bytes:
    presentation = Presentation()
    slide = presentation.slides.add_slide(presentation.slide_layouts[1])
    slide.shapes.title.text = "经营简报"
    slide.placeholders[1].text = "两个项目进入延期关注。"
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


def _pdf() -> bytes:
    writer = PdfWriter()
    page = writer.add_blank_page(width=595, height=842)
    font = DictionaryObject(
        {
            NameObject("/Type"): NameObject("/Font"),
            NameObject("/Subtype"): NameObject("/Type1"),
            NameObject("/BaseFont"): NameObject("/Helvetica"),
        }
    )
    page[NameObject("/Resources")] = DictionaryObject(
        {NameObject("/Font"): DictionaryObject({NameObject("/F1"): writer._add_object(font)})}
    )
    stream = DecodedStreamObject()
    stream.set_data(b"BT /F1 12 Tf 72 720 Td (Executive business brief) Tj ET")
    page[NameObject("/Contents")] = writer._add_object(stream)
    output = io.BytesIO()
    writer.write(output)
    return output.getvalue()


def test_supported_office_and_pdf_parsers_preserve_locators() -> None:
    cases = [
        ("brief.docx", _docx(), "python-docx"),
        ("collections.xlsx", _xlsx(), "openpyxl"),
        ("brief.pptx", _pptx(), "python-pptx"),
        ("appendix.pdf", _pdf(), "pypdf"),
    ]
    for name, content, parser_name in cases:
        blocks, page_count, actual_parser = parse_file(name, content)
        assert actual_parser == parser_name
        assert page_count >= 1
        assert blocks
        assert all(block.locator for block in blocks)


def test_office_archive_rejects_excessive_expansion(monkeypatch) -> None:
    payload = io.BytesIO()
    with zipfile.ZipFile(payload, "w", compression=zipfile.ZIP_DEFLATED) as archive:
        archive.writestr("word/document.xml", b"0" * 2048)
    monkeypatch.setattr(
        "executive_ai_worker.file_extraction.MAX_OFFICE_UNCOMPRESSED_BYTES",
        1024,
    )
    try:
        parse_file("oversized.docx", payload.getvalue())
    except FileExtractionPermanentError as exc:
        assert exc.code == "document_resource_limit"
    else:  # pragma: no cover - the archive must be rejected before parsing.
        raise AssertionError("excessively expanded Office archive was accepted")


def test_extracted_text_limit_is_enforced(monkeypatch) -> None:
    monkeypatch.setattr("executive_ai_worker.file_extraction.MAX_EXTRACTED_CHARACTERS", 5)
    try:
        parse_file("brief.docx", _docx())
    except FileExtractionPermanentError as exc:
        assert exc.code == "document_resource_limit"
    else:  # pragma: no cover - extracted text must be bounded.
        raise AssertionError("oversized extracted text was accepted")


def test_chunk_count_limit_applies_to_small_blocks(monkeypatch) -> None:
    monkeypatch.setattr("executive_ai_worker.file_extraction.MAX_FILE_CHUNKS", 1)
    try:
        chunk_blocks(
            [
                ExtractedBlock("one", {"type": "paragraph", "paragraph": 1}),
                ExtractedBlock("two", {"type": "paragraph", "paragraph": 2}),
            ]
        )
    except FileExtractionPermanentError as exc:
        assert exc.code == "document_resource_limit"
    else:  # pragma: no cover - every chunk path must enforce the same cap.
        raise AssertionError("excessive small chunks were accepted")


def test_runtime_embedding_loader_is_strictly_offline(monkeypatch, tmp_path: Path) -> None:
    captured: dict[str, object] = {}

    class FakeTextEmbedding:
        def __init__(self, **kwargs: object) -> None:
            captured.update(kwargs)

    monkeypatch.setattr("fastembed.TextEmbedding", FakeTextEmbedding)
    _embedding_model.cache_clear()
    _embedding_model("test/model", str(tmp_path))
    assert captured["cache_dir"] == str(tmp_path)
    assert captured["local_files_only"] is True
    _embedding_model.cache_clear()


def test_hybrid_file_retrieval_terms_and_rank_fusion() -> None:
    class Chunk:
        def __init__(self) -> None:
            self.id = uuid.uuid4()

    vector_first = (Chunk(), object())
    shared = (Chunk(), object())
    keyword_first = (Chunk(), object())
    terms = _query_terms("请核对合同号 HT-2026-001 和华东事业部回款")
    ranked = _reciprocal_rank_fusion(
        [vector_first, shared],
        [shared, keyword_first],
        limit=3,
    )

    assert "ht-2026-001" in terms
    assert "华东事业部" in terms
    assert "回款" in terms
    assert ranked[0][0].id == shared[0].id


def test_daily_schedule_uses_configured_timezone_and_next_window() -> None:
    after = datetime(2026, 7, 26, 17, 59, tzinfo=UTC)
    next_run = next_cron_time("0 2 * * *", "Asia/Shanghai", after)
    assert next_run == datetime(2026, 7, 26, 18, 0, tzinfo=UTC)


def _embedding_bundle(tmp_path: Path) -> tuple[Path, str]:
    source = tmp_path / "source" / "fast-test-model"
    source.mkdir(parents=True)
    (source / "model_optimized.onnx").write_bytes(b"verified-model")
    (source / "tokenizer.json").write_text("{}")
    archive = tmp_path / "model.tar.gz"
    with tarfile.open(archive, "w:gz") as bundle:
        bundle.add(source, arcname="fast-test-model")
    return archive, hashlib.sha256(archive.read_bytes()).hexdigest()


def test_embedding_preload_is_integrity_checked_and_idempotent(
    tmp_path: Path,
    monkeypatch,
) -> None:
    archive, sha256 = _embedding_bundle(tmp_path)
    artifact = EmbeddingArtifact(
        model_name="test/model",
        url="https://artifacts.test/model.tar.gz",
        sha256=sha256,
        archive_root="fast-test-model",
        required_file="model_optimized.onnx",
    )
    cache = tmp_path / "cache"
    monkeypatch.setattr(
        "executive_ai_worker.embedding_cache._download_with_resume",
        lambda _url, target: (
            target.parent.mkdir(parents=True, exist_ok=True),
            target.write_bytes(archive.read_bytes()),
        ),
    )
    first = preload_artifact(artifact, cache)
    archive.unlink()
    second = preload_artifact(artifact, cache)
    assert first == second
    assert (second / "model_optimized.onnx").read_bytes() == b"verified-model"
    assert (second / ".artifact-sha256").read_text().strip() == sha256


def test_embedding_preload_rejects_unapproved_artifact_hash(tmp_path: Path, monkeypatch) -> None:
    archive, _ = _embedding_bundle(tmp_path)
    artifact = EmbeddingArtifact(
        model_name="test/model",
        url="https://artifacts.test/model.tar.gz",
        sha256="0" * 64,
        archive_root="fast-test-model",
        required_file="model_optimized.onnx",
    )
    monkeypatch.setattr(
        "executive_ai_worker.embedding_cache._download_with_resume",
        lambda _url, target: (
            target.parent.mkdir(parents=True, exist_ok=True),
            target.write_bytes(archive.read_bytes()),
        ),
    )
    monkeypatch.setattr("executive_ai_worker.embedding_cache.time.sleep", lambda _: None)
    try:
        preload_artifact(artifact, tmp_path / "cache")
    except RuntimeError as exc:
        assert "download failed" in str(exc)
    else:  # pragma: no cover - a hash mismatch must always be terminal.
        raise AssertionError("unapproved embedding artifact was accepted")
